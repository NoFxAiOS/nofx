#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Script para crear presentación PowerPoint de la Auditoría PROAGUA 2024
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle):
    """Crea la diapositiva de título"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    # Estilo del título
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

def create_content_slide(prs, title, content_items):
    """Crea diapositiva de contenido con viñetas"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    # Título
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Contenido
    content = slide.placeholders[1]
    text_frame = content.text_frame
    text_frame.clear()

    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)

def create_table_slide(prs, title, headers, rows):
    """Crea diapositiva con tabla"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Título
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Tabla
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(0.5)

    table = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table

    # Encabezados
    for i, header in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.bold = True

    # Datos
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.rows[i + 1].cells[j]
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(12)

def main():
    # Crear presentación
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. Diapositiva de título
    create_title_slide(prs,
                      "AUDITORÍA PROAGUA 2024",
                      "Gobierno del Estado de México\nAuditoría de Cumplimiento 2024-A-15000-19-0856-2025")

    # 2. Información General
    create_content_slide(prs, "1. Información General", [
        "Auditoría No. 856 - Modalidad Presencial",
        "Programa de Agua Potable, Drenaje y Tratamiento (PROAGUA)",
        "Tipo: Auditoría de Cumplimiento 2024",
        "Estado: Auditoría Superior de la Federación"
    ])

    # 3. Recursos y Alcance
    create_content_slide(prs, "2. Recursos y Alcance", [
        "Universo Total del Programa: $2,606,615.1 miles de pesos",
        "Recursos Transferidos al Estado de México: $197,904.0 miles",
        "Muestra Auditada: $197,904.0 miles (100%)",
        "Rendimientos Financieros: $13,145.9 miles",
        "Recursos Disponibles: $211,058.4 miles de pesos"
    ])

    # 4. Destino de los Recursos
    create_content_slide(prs, "3. Destino de los Recursos", [
        "Total de Proyectos: 53 proyectos",
        "Importe Pagado (al 31/julio/2025): $134,099.3 miles (63.6%)",
        "Componentes principales:",
        "  • Agua Potable Rehabilitado: 29 proyectos (46.9%)",
        "  • Mejoramiento de eficiencia: 14 proyectos (8.0%)",
        "  • Supervisión, Atención Social y Fortalecimiento"
    ])

    # 5. Tabla de Distribución por Componente
    create_table_slide(prs, "Distribución por Componente",
                      ["Componente", "Proyectos", "Importe (miles $)", "%"],
                      [
                          ["Agua Potable - Nuevo", "3", "6,472.4", "3.1%"],
                          ["Agua Potable - Rehabilitado", "29", "98,922.2", "46.9%"],
                          ["Mejoramiento de eficiencia", "14", "16,786.6", "8.0%"],
                          ["Alcantarillado", "3", "3,511.9", "1.7%"],
                          ["Desinfección", "1", "4,500.0", "2.1%"],
                          ["Supervisión y Otros", "3", "3,906.2", "1.9%"]
                      ])

    # 6. Dependencias Ejecutoras
    create_content_slide(prs, "4. Dependencias Ejecutoras", [
        "CAEM (Comisión del Agua): $122,408.2 miles",
        "ODAPAS Chalco: $4,936.9 miles",
        "Aguas de Huixquilucan: $4,705.8 miles",
        "OPDAPAS San Mateo Atenco: $1,950.0 miles",
        "APAST Tultitlán: $540.0 miles",
        "ODAPAZ Zumpango: $2,000.0 miles"
    ])

    # 7. Hallazgos Positivos
    create_content_slide(prs, "5. Hallazgos Positivos ✓", [
        "Transparencia en Adjudicación:",
        "  • 25 contratos adjudicados conforme a normativa",
        "  • 8 por licitación pública ($156,840.6 miles)",
        "  • 12 por invitación a 3 personas ($76,478.6 miles)",
        "",
        "Verificación Física Exitosa:",
        "  • 16 obras concluidas y operando adecuadamente",
        "  • Registros contables actualizados y controlados"
    ])

    # 8. Observaciones Críticas
    create_content_slide(prs, "6. Observaciones Críticas ⚠", [
        "1. Recursos No Reintegrados: $23,676.4 miles",
        "   Pendiente: $13,165.4 miles de pesos",
        "",
        "2. Pagos Sin Documentación: $652.1 miles",
        "   Empresa EFECTIVALE - Combustible y alimentos",
        "",
        "3. Obra Sin Propiedad del Predio: $1,467.7 miles",
        "",
        "4. Anticipo No Amortizado: $2,662.5 miles"
    ])

    # 9. Incumplimientos Administrativos
    create_content_slide(prs, "7. Incumplimientos Administrativos", [
        "Gestión Bancaria:",
        "  ✗ Cuenta bancaria no exclusiva para recursos",
        "  ✗ Transferencias fuera de plazo (>10 días)",
        "  ✗ 4 ejecutores sin cuentas exclusivas",
        "",
        "Documentación:",
        "  ✗ 6 contratos sin opinión fiscal del SAT",
        "  ✗ 5 obras sin título de concesión CONAGUA",
        "  ✗ 6 obras sin cierre administrativo completo"
    ])

    # 10. Tabla de Daño a Hacienda Pública
    create_table_slide(prs, "8. Daño a la Hacienda Pública",
                      ["No.", "Concepto", "Monto (pesos)", "Estado"],
                      [
                          ["001", "Pagos sin devengación", "$652,098.79", "Pendiente"],
                          ["002", "Recursos no reintegrados", "$13,165,438.95", "Parcial"],
                          ["003", "Obra sin propiedad", "$1,467,724.80", "Pendiente"],
                          ["004", "Anticipo no amortizado", "$2,662,537.23", "Pendiente"],
                          ["", "TOTAL", "$17,947,799.77", ""]
                      ])

    # 11. Transparencia y Rendición de Cuentas
    create_content_slide(prs, "9. Transparencia y Rendición de Cuentas", [
        "Cumplimientos ✓:",
        "  • Información reportada a SHCP mediante SRFT (CAEM)",
        "  • Difusión en página de internet",
        "  • Leyenda normativa en publicidad",
        "",
        "Incumplimientos ✗:",
        "  • 5 dependencias ejecutoras NO reportaron:",
        "    ODAPAS Chalco, Aguas de Huixquilucan,",
        "    OPDAPAS San Mateo, APAST Tultitlán, ODAPAZ Zumpango"
    ])

    # 12. Indicadores de Desempeño
    create_content_slide(prs, "10. Indicadores de Desempeño", [
        "Cobertura de Agua Potable:",
        "  • Meta normativa: ≥97.5%",
        "  • Cobertura alcanzada: 97.1%",
        "  • Estado: ⚠ Por debajo de la meta",
        "",
        "Límites de Gasto Especial:",
        "  ✓ Desinfección: $4,500.0 miles (dentro de límite)",
        "  ✓ Fortalecimiento: $972.9 miles (dentro de límite)"
    ])

    # 13. Acciones Promovidas
    create_content_slide(prs, "11. Acciones Promovidas", [
        "Pliegos de Observaciones: 4",
        "  • Total Observado: $28,540,922.05",
        "  • Recuperado durante auditoría: $10,593,122.28",
        "  • Pendiente: $17,947,799.77",
        "",
        "Procedimientos Administrativos: 19 expedientes",
        "  • OIC Secretaría de Finanzas",
        "  • OIC CAEM",
        "  • OIC Organismos Municipales (6)"
    ])

    # 14. Conclusiones - Fortalezas
    create_content_slide(prs, "12. Conclusiones - Fortalezas", [
        "Dictamen: GESTIÓN RAZONABLE CON ÁREAS DE OPORTUNIDAD",
        "",
        "✓ 100% de recursos auditados",
        "✓ Obras físicamente verificadas y en operación",
        "✓ Procesos de adjudicación transparentes",
        "✓ Registro contable adecuado",
        "✓ Recuperación de $10,593,122.28 durante auditoría"
    ])

    # 15. Conclusiones - Debilidades
    create_content_slide(prs, "13. Conclusiones - Debilidades", [
        "✗ Control interno deficiente en gestión bancaria",
        "✗ Falta de documentación soporte en pagos",
        "✗ Incumplimiento de contraparte estatal:",
        "   Comprometida: $230,850.6 miles",
        "   Realizada: $174,552.6 miles",
        "   Faltante: $56,298.0 miles",
        "✗ Reintegros fuera de plazo",
        "✗ Obras sin cierre administrativo",
        "✗ Transparencia incompleta (ejecutores municipales)"
    ])

    # 16. Recomendaciones Clave
    create_content_slide(prs, "14. Recomendaciones Clave", [
        "1. Gestión Financiera:",
        "   • Cuentas bancarias exclusivas por programa",
        "   • Transferencias dentro de plazos normativos",
        "",
        "2. Control Documental:",
        "   • Verificación de opinión fiscal antes de contratar",
        "   • Cierre administrativo oportuno de obras",
        "",
        "3. Cumplimiento Normativo:",
        "   • Garantizar contraparte estatal comprometida",
        "   • Reportes completos de todos los ejecutores"
    ])

    # 17. Top 5 Proyectos
    create_table_slide(prs, "15. Principales Proyectos Revisados",
                      ["No.", "Proyecto", "Monto (miles $)"],
                      [
                          ["1", "Rehabilitación Jocotitlán", "$14,830.1"],
                          ["2", "Ampliación Taborda", "$6,956.2"],
                          ["3", "Equipamiento Tepexpan", "$6,338.7"],
                          ["4", "Rehabilitación Jaltenco", "$5,677.9"],
                          ["5", "Rehabilitación Santiago Huitlapaltepec", "$5,664.6"]
                      ])

    # 18. Diapositiva final
    create_content_slide(prs, "Datos de Contacto", [
        "Auditoría Superior de la Federación",
        "",
        "Director de Área: Ing. Daniel Palacios Téllez",
        "Director General: Lic. Marciano Cruz Ángeles",
        "",
        "Fecha de Informe: 2025",
        "Tipo: Auditoría de Cumplimiento - Cuenta Pública 2024",
        "",
        "Documento base: Informe Individual del Resultado de",
        "la Fiscalización Superior - Auditoría No. 856"
    ])

    # Guardar presentación
    output_file = "/vercel/sandbox/Auditoria_PROAGUA_2024.pptx"
    prs.save(output_file)
    print(f"✓ Presentación creada exitosamente: {output_file}")
    print(f"✓ Total de diapositivas: {len(prs.slides)}")

if __name__ == "__main__":
    main()
